"""Knowledge Store - Persistent RAG storage for user documents and research."""

import os
import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional

import chromadb
from chromadb.utils import embedding_functions

logger = logging.getLogger(__name__)

KNOWLEDGE_BASE_DIR = Path("data/knowledge")

_embedding_function_cache = None
_knowledge_store_cache: Dict[str, "KnowledgeStore"] = {}


def get_embedding_function():
    """Singleton for SentenceTransformer embedding function - loads once, reused forever."""
    global _embedding_function_cache
    if _embedding_function_cache is None:
        logger.info("[KNOWLEDGE] Loading SentenceTransformer model (one-time)...")
        _embedding_function_cache = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name="all-MiniLM-L6-v2"
        )
        logger.info("[KNOWLEDGE] SentenceTransformer model loaded and cached.")
    return _embedding_function_cache


def get_knowledge_store(session_id: str) -> "KnowledgeStore":
    """Get or create KnowledgeStore for session (cached by session_id)."""
    global _knowledge_store_cache
    if session_id not in _knowledge_store_cache:
        _knowledge_store_cache[session_id] = KnowledgeStore(session_id)
    return _knowledge_store_cache[session_id]


def invalidate_knowledge_store_cache(session_id: str = None):
    """Invalidate cached KnowledgeStore instance(s) to force fresh reads."""
    global _knowledge_store_cache
    if session_id:
        if session_id in _knowledge_store_cache:
            del _knowledge_store_cache[session_id]
            logger.info(f"[KNOWLEDGE] Invalidated cache for session {session_id}")
    else:
        _knowledge_store_cache = {}
        logger.info("[KNOWLEDGE] Invalidated all cached stores")


class KnowledgeStore:
    """Per-session persistent vector store for documents and research."""

    def __init__(self, session_id: str):
        self.session_id = session_id
        self.storage_path = KNOWLEDGE_BASE_DIR / session_id
        self.storage_path.mkdir(parents=True, exist_ok=True)

        self.client = chromadb.PersistentClient(path=str(self.storage_path))
        self.ef = get_embedding_function()
        self.r2_service = self._get_r2_service()

        try:
            self.collection = self.client.get_or_create_collection(name="knowledge", embedding_function=self.ef)
            self.memory_collection = self.client.get_or_create_collection(
                name="conversation_memory", embedding_function=self.ef
            )
            logger.info(f"[KNOWLEDGE] Initialized store for session {session_id}")

            if self.collection.count() == 0 and self.r2_service:
                self._restore_from_r2()
        except Exception as e:
            logger.error(f"[KNOWLEDGE] Failed to create collection: {e}")
            self.collection = None
            self.memory_collection = None

    def _get_r2_service(self):
        try:
            from src.storage.cloud_storage import get_cloud_storage

            return get_cloud_storage()
        except Exception:
            return None

    def _restore_from_r2(self):
        if not self.r2_service:
            return
        try:
            prefix = f"{self.session_id}/knowledge/"
            blobs = self.r2_service.list_blobs(prefix=prefix)
            if not blobs:
                return
            logger.info(f"[KNOWLEDGE] Restoring {len(blobs)} documents from R2")
            for blob in blobs:
                try:
                    local_tmp = self.storage_path / "tmp_restore.json"
                    self.r2_service.download_file(blob["key"], local_tmp)
                    with open(local_tmp, "r", encoding="utf-8") as f:
                        doc_data = json.load(f)
                    local_tmp.unlink()
                    self.collection.add(
                        documents=[doc_data["content"]], metadatas=[doc_data["metadata"]], ids=[doc_data["id"]]
                    )
                except Exception as e:
                    logger.warning(f"[KNOWLEDGE] Failed to restore {blob['key']}: {e}")
        except Exception as e:
            logger.warning(f"[KNOWLEDGE] R2 restore failed: {e}")

    def _sync_to_r2(self, doc_id: str, content: str, metadata: Dict):
        if not self.r2_service:
            return
        try:
            doc_data = {"id": doc_id, "content": content, "metadata": metadata}
            local_tmp = self.storage_path / f"{doc_id}.json"
            with open(local_tmp, "w", encoding="utf-8") as f:
                json.dump(doc_data, f)
            blob_path = f"{self.session_id}/knowledge/{doc_id}.json"
            self.r2_service.upload_file(local_tmp, blob_path)
            local_tmp.unlink()
            logger.info(f"[KNOWLEDGE] Synced {doc_id} to R2")
        except Exception as e:
            logger.warning(f"[KNOWLEDGE] R2 sync failed for {doc_id}: {e}")

    def _delete_from_r2(self, doc_id: str):
        if not self.r2_service:
            return
        try:
            blob_path = f"{self.session_id}/knowledge/{doc_id}.json"
            self.r2_service.delete_blob(blob_path)
            logger.info(f"[KNOWLEDGE] Deleted {doc_id} from R2")
        except Exception:
            pass

    def _generate_id(self, content: str) -> str:
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def add_document(
        self, content: str, source: str = "user", source_name: str = "", metadata: Optional[Dict] = None
    ) -> str:
        if not self.collection:
            return ""

        doc_id = f"doc_{self._generate_id(content)}"
        doc_metadata = {
            "source": source,
            "source_name": source_name or "Unknown",
            "added_at": datetime.now().isoformat(),
            "content_preview": content[:200],
            "inject_to_context": source not in ["dataset", "dataset_profile"],
            **(metadata or {}),
        }

        try:
            self.collection.add(documents=[content], metadatas=[doc_metadata], ids=[doc_id])
            self._sync_to_r2(doc_id, content, doc_metadata)
            logger.info(f"[KNOWLEDGE] Added document {doc_id} from {source}")
            return doc_id
        except Exception as e:
            logger.error(f"[KNOWLEDGE] Failed to add document: {e}")
            return ""

    def add_file(self, file_path: str, source_type: str = "user_upload") -> str:
        """Add a file to knowledge store by extracting its text content."""
        path = Path(file_path)
        if not path.exists():
            logger.error(f"[KNOWLEDGE] File not found: {file_path}")
            return ""

        try:
            content = self._extract_file_content(path)
            if not content:
                return ""

            return self.add_document(
                content=content,
                source=source_type,
                source_name=path.name,
                metadata={"file_path": str(file_path), "file_type": path.suffix},
            )
        except Exception as e:
            logger.error(f"[KNOWLEDGE] Failed to process file {file_path}: {e}")
            return ""

    def _extract_file_content(self, path: Path) -> str:
        """Extract text content from various file types."""
        suffix = path.suffix.lower()

        # Plain text formats
        if suffix in [".txt", ".md", ".rst", ".org", ".py", ".json", ".log"]:
            return path.read_text(encoding="utf-8", errors="ignore")

        # CSV/TSV
        if suffix in [".csv", ".tsv"]:
            try:
                import pandas as pd

                sep = "\t" if suffix == ".tsv" else ","
                df = pd.read_csv(path, sep=sep, nrows=1000)
                return df.to_string()
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] CSV/TSV extraction failed: {e}")
                return path.read_text(encoding="utf-8", errors="ignore")

        # PDF
        if suffix == ".pdf":
            try:
                import fitz

                doc = fitz.open(path)
                text = "\n".join(page.get_text() for page in doc)
                doc.close()
                return text
            except ImportError:
                logger.warning("[KNOWLEDGE] PyMuPDF not installed")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] PDF extraction failed: {e}")
            return ""

        # Microsoft Word (.docx)
        if suffix == ".docx":
            try:
                from docx import Document

                doc = Document(path)
                return "\n".join(para.text for para in doc.paragraphs)
            except ImportError:
                logger.warning("[KNOWLEDGE] python-docx not installed")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] DOCX extraction failed: {e}")
            return ""

        # Microsoft PowerPoint (.pptx)
        if suffix == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(path)
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                return "\n".join(text_parts)
            except ImportError:
                logger.warning("[KNOWLEDGE] python-pptx not installed")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] PPTX extraction failed: {e}")
            return ""

        # Microsoft Excel (.xlsx, .xls)
        if suffix in [".xlsx", ".xls"]:
            try:
                import pandas as pd

                df = pd.read_excel(path, nrows=1000)
                return df.to_string()
            except ImportError:
                logger.warning("[KNOWLEDGE] openpyxl not installed")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] Excel extraction failed: {e}")
            return ""

        # HTML/XML
        if suffix in [".html", ".htm", ".xml"]:
            try:
                from bs4 import BeautifulSoup

                content = path.read_text(encoding="utf-8", errors="ignore")
                soup = BeautifulSoup(content, "html.parser")
                for tag in soup(["script", "style", "nav", "footer", "header"]):
                    tag.decompose()
                return soup.get_text(separator="\n", strip=True)
            except ImportError:
                logger.warning("[KNOWLEDGE] BeautifulSoup not installed")
                return path.read_text(encoding="utf-8", errors="ignore")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] HTML/XML extraction failed: {e}")
            return ""

        # Email (.eml)
        if suffix == ".eml":
            try:
                import email
                from email import policy

                with open(path, "rb") as f:
                    msg = email.message_from_binary_file(f, policy=policy.default)
                body = msg.get_body(preferencelist=("plain", "html"))
                if body:
                    content = body.get_content()
                    if "<html" in content.lower():
                        from bs4 import BeautifulSoup

                        return BeautifulSoup(content, "html.parser").get_text()
                    return content
                return f"Subject: {msg['subject']}\nFrom: {msg['from']}\nTo: {msg['to']}"
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] EML extraction failed: {e}")
            return ""

        # Open Document Text (.odt)
        if suffix == ".odt":
            try:
                from odf import text as odf_text
                from odf.opendocument import load

                doc = load(path)
                paragraphs = doc.getElementsByType(odf_text.P)
                return "\n".join(str(p) for p in paragraphs)
            except ImportError:
                logger.warning("[KNOWLEDGE] odfpy not installed")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] ODT extraction failed: {e}")
            return ""

        # RTF
        if suffix == ".rtf":
            try:
                from striprtf.striprtf import rtf_to_text

                content = path.read_text(encoding="utf-8", errors="ignore")
                return rtf_to_text(content)
            except ImportError:
                logger.warning("[KNOWLEDGE] striprtf not installed")
            except Exception as e:
                logger.warning(f"[KNOWLEDGE] RTF extraction failed: {e}")
            return ""

        logger.warning(f"[KNOWLEDGE] Unsupported file type: {suffix}")
        return ""

    def query(self, text: str, k: int = 5) -> List[Dict[str, Any]]:
        """Query knowledge store for relevant documents."""
        if not self.collection:
            return []

        try:
            results = self.collection.query(query_texts=[text], n_results=k)

            items = []
            if results["documents"] and results["documents"][0]:
                for i, doc in enumerate(results["documents"][0]):
                    items.append(
                        {
                            "content": doc,
                            "metadata": results["metadatas"][0][i] if results["metadatas"] else {},
                            "id": results["ids"][0][i] if results["ids"] else "",
                        }
                    )
            return items

        except Exception as e:
            logger.error(f"[KNOWLEDGE] Query failed: {e}")
            return []

    def _refresh_collection(self):
        """Refresh collection reference to get latest data from SQLite."""
        try:
            self.collection = self.client.get_or_create_collection(name="knowledge", embedding_function=self.ef)
        except Exception as e:
            logger.warning(f"[KNOWLEDGE] Failed to refresh collection: {e}")

    def list_items(self) -> List[Dict[str, Any]]:
        """List all items in knowledge store."""
        if not self.collection:
            return []

        try:
            self._refresh_collection()
            count = self.collection.count()
            logger.info(f"[KNOWLEDGE] list_items: collection has {count} documents")

            if count == 0:
                return []

            results = self.collection.get(include=["metadatas"])
            items = []
            for i, doc_id in enumerate(results["ids"]):
                meta = results["metadatas"][i] if results["metadatas"] else {}
                items.append(
                    {
                        "id": doc_id,
                        "source": meta.get("source", "unknown"),
                        "source_name": meta.get("source_name", "Unknown"),
                        "added_at": meta.get("added_at", ""),
                        "content_preview": meta.get("content_preview", ""),
                        "inject_to_context": meta.get("inject_to_context", False),
                    }
                )
            logger.info(f"[KNOWLEDGE] list_items: returning {len(items)} items")
            return items
        except Exception as e:
            logger.error(f"[KNOWLEDGE] List failed: {e}")
            return []

    def delete_item(self, doc_id: str) -> bool:
        if not self.collection:
            return False

        try:
            self.collection.delete(ids=[doc_id])
            self._delete_from_r2(doc_id)
            logger.info(f"[KNOWLEDGE] Deleted {doc_id}")
            return True
        except Exception as e:
            logger.error(f"[KNOWLEDGE] Delete failed: {e}")
            return False

    def count(self) -> int:
        if not self.collection:
            return 0
        return self.collection.count()

    def get_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        if not self.collection:
            return None
        try:
            results = self.collection.get(ids=[doc_id], include=["documents", "metadatas"])
            if not results["ids"]:
                return None
            return {
                "id": doc_id,
                "content": results["documents"][0] if results["documents"] else "",
                "metadata": results["metadatas"][0] if results["metadatas"] else {},
            }
        except Exception as e:
            logger.error(f"[KNOWLEDGE] Get document failed: {e}")
            return None

    def list_injectable_items(self) -> List[Dict[str, Any]]:
        """List items where inject_to_context is True (for Brain context injection)."""
        if not self.collection:
            return []

        try:
            results = self.collection.get(include=["documents", "metadatas"])
            items = []
            for i, doc_id in enumerate(results.get("ids", [])):
                meta = results["metadatas"][i] if results.get("metadatas") else {}
                if meta.get("inject_to_context", False):
                    items.append(
                        {
                            "id": doc_id,
                            "source_name": meta.get("source_name", "Unknown"),
                            "content": results["documents"][i] if results.get("documents") else "",
                        }
                    )
            return items
        except Exception as e:
            logger.error(f"[KNOWLEDGE] List injectable failed: {e}")
            return []

    def update_item_metadata(self, doc_id: str, updates: Dict[str, Any]) -> bool:
        """Update metadata for an existing item."""
        if not self.collection:
            return False

        try:
            existing = self.collection.get(ids=[doc_id], include=["metadatas"])
            if not existing["ids"]:
                return False

            current_meta = existing["metadatas"][0] if existing.get("metadatas") else {}
            current_meta.update(updates)
            self.collection.update(ids=[doc_id], metadatas=[current_meta])
            logger.info(f"[KNOWLEDGE] Updated metadata for {doc_id}")
            return True
        except Exception as e:
            logger.error(f"[KNOWLEDGE] Update metadata failed: {e}")
            return False

    def add_conversation_turn(self, user_message: str, ai_response: str, turn_index: int = None) -> str:
        if not self.memory_collection:
            return ""

        try:
            combined_text = f"User: {user_message}\n\nAssistant: {ai_response[:1500]}"
            turn_id = f"turn_{self._generate_id(combined_text)}_{datetime.now().strftime('%H%M%S')}"

            metadata = {
                "user_message": user_message[:500],
                "response_preview": ai_response[:300],
                "timestamp": datetime.now().isoformat(),
                "turn_index": turn_index or self.memory_collection.count(),
            }

            self.memory_collection.add(documents=[combined_text], metadatas=[metadata], ids=[turn_id])
            logger.info(f"[MEMORY] Stored conversation turn {turn_id}")
            return turn_id
        except Exception as e:
            logger.error(f"[MEMORY] Failed to store turn: {e}")
            return ""

    def get_relevant_history(self, query: str, k: int = 5, min_score: float = 0.3) -> List[Dict[str, Any]]:
        if not self.memory_collection or self.memory_collection.count() == 0:
            return []

        try:
            results = self.memory_collection.query(
                query_texts=[query],
                n_results=min(k, self.memory_collection.count()),
                include=["documents", "metadatas", "distances"],
            )

            items = []
            if results["documents"] and results["documents"][0]:
                for i, doc in enumerate(results["documents"][0]):
                    distance = results["distances"][0][i] if results.get("distances") else 1.0
                    relevance_score = 1.0 - distance

                    if relevance_score >= min_score:
                        items.append(
                            {
                                "content": doc,
                                "metadata": results["metadatas"][0][i] if results.get("metadatas") else {},
                                "relevance": relevance_score,
                            }
                        )

            items.sort(key=lambda x: x["metadata"].get("turn_index", 0))
            return items
        except Exception as e:
            logger.error(f"[MEMORY] Query failed: {e}")
            return []

    def get_recent_history(self, k: int = 10) -> List[Dict[str, Any]]:
        if not self.memory_collection or self.memory_collection.count() == 0:
            return []

        try:
            results = self.memory_collection.get(include=["documents", "metadatas"])
            items = []
            for i, doc_id in enumerate(results.get("ids", [])):
                items.append(
                    {
                        "id": doc_id,
                        "content": results["documents"][i] if results.get("documents") else "",
                        "metadata": results["metadatas"][i] if results.get("metadatas") else {},
                    }
                )

            items.sort(key=lambda x: x["metadata"].get("turn_index", 0), reverse=True)
            return items[:k]
        except Exception as e:
            logger.error(f"[MEMORY] Get recent failed: {e}")
            return []

    def clear_conversation_memory(self) -> bool:
        if not self.memory_collection:
            return False

        try:
            all_ids = self.memory_collection.get()["ids"]
            if all_ids:
                self.memory_collection.delete(ids=all_ids)
            logger.info(f"[MEMORY] Cleared conversation memory for {self.session_id}")
            return True
        except Exception as e:
            logger.error(f"[MEMORY] Clear failed: {e}")
            return False

    def memory_count(self) -> int:
        if not self.memory_collection:
            return 0
        return self.memory_collection.count()
