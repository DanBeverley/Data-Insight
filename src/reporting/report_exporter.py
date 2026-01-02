import io
import re
from html import unescape
from typing import Dict, List, Any, Optional
from datetime import datetime
from pathlib import Path


def html_to_text(html_content: str) -> str:
    text = re.sub(r"<script[^>]*>.*?</script>", "", html_content, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<h1[^>]*>(.*?)</h1>", r"\n\n=== \1 ===\n\n", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<h2[^>]*>(.*?)</h2>", r"\n\n## \1\n\n", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<h3[^>]*>(.*?)</h3>", r"\n\n### \1\n\n", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<li[^>]*>(.*?)</li>", r"• \1\n", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<p[^>]*>(.*?)</p>", r"\1\n\n", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<strong[^>]*>(.*?)</strong>", r"**\1**", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<b[^>]*>(.*?)</b>", r"**\1**", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<em[^>]*>(.*?)</em>", r"*\1*", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(
        r'<iframe[^>]*data-filename="([^"]+)"[^>]*>.*?</iframe>',
        r"\n[Chart: \1]\n",
        text,
        flags=re.DOTALL | re.IGNORECASE,
    )
    text = re.sub(
        r'<div[^>]*data-filename="([^"]+)"[^>]*>.*?</div>', r"\n[Chart: \1]\n", text, flags=re.DOTALL | re.IGNORECASE
    )
    text = re.sub(r"<[^>]+>", "", text)
    text = unescape(text)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    return text.strip()


def html_to_reportlab(html_content: str) -> str:
    text = re.sub(r"<script[^>]*>.*?</script>", "", html_content, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(
        r"<h1[^>]*>(.*?)</h1>", r'<font size="18"><b>\1</b></font><br/><br/>', text, flags=re.DOTALL | re.IGNORECASE
    )
    text = re.sub(
        r"<h2[^>]*>(.*?)</h2>", r'<font size="14"><b>\1</b></font><br/>', text, flags=re.DOTALL | re.IGNORECASE
    )
    text = re.sub(
        r"<h3[^>]*>(.*?)</h3>", r'<font size="12"><b>\1</b></font><br/>', text, flags=re.DOTALL | re.IGNORECASE
    )
    text = re.sub(r"<li[^>]*>(.*?)</li>", r"• \1<br/>", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<p[^>]*>(.*?)</p>", r"\1<br/><br/>", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<br\s*/?>", "<br/>", text, flags=re.IGNORECASE)
    text = re.sub(r"<strong[^>]*>(.*?)</strong>", r"<b>\1</b>", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<em[^>]*>(.*?)</em>", r"<i>\1</i>", text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(
        r'<iframe[^>]*src="[^"]*?([^/"]+\.html)"[^>]*>.*?</iframe>',
        r"<br/>[Chart: \1]<br/>",
        text,
        flags=re.DOTALL | re.IGNORECASE,
    )
    text = re.sub(
        r'<div[^>]*data-filename="([^"]+)"[^>]*>.*?</div>',
        r"<br/>[Chart: \1]<br/>",
        text,
        flags=re.DOTALL | re.IGNORECASE,
    )
    text = re.sub(r"<ul[^>]*>|</ul>|<ol[^>]*>|</ol>", "", text, flags=re.IGNORECASE)
    text = re.sub(r"<div[^>]*>|</div>", "", text, flags=re.IGNORECASE)
    text = re.sub(r"<span[^>]*>|</span>", "", text, flags=re.IGNORECASE)
    text = unescape(text)
    text = re.sub(r"(<br/>){3,}", "<br/><br/>", text)
    return text.strip()


def make_standalone_html(html_content: str) -> str:
    artifact_pattern = r'<iframe[^>]*src="([^"]*?)([^/"]+\.html)"[^>]*>.*?</iframe>'
    matches = re.findall(artifact_pattern, html_content, flags=re.DOTALL | re.IGNORECASE)

    plotly_cdn = '<script src="https://cdn.plot.ly/plotly-2.27.0.min.js"></script>'
    cdn_added = False

    for prefix, filename in matches:
        for base_path in [Path("static/plots"), Path("data/reports")]:
            artifact_path = base_path / filename
            if artifact_path.exists():
                try:
                    artifact_html = artifact_path.read_text(encoding="utf-8")
                    replacement = f'<div class="embedded-chart" style="width:100%;height:500px;overflow:auto;border:1px solid #334155;margin:10px 0;">{artifact_html}</div>'
                    original_src = f"{prefix}{filename}"
                    html_content = re.sub(
                        rf'<iframe[^>]*src="{re.escape(original_src)}"[^>]*>.*?</iframe>',
                        replacement,
                        html_content,
                        flags=re.DOTALL | re.IGNORECASE,
                    )
                    cdn_added = True
                except Exception:
                    pass
                break

    if cdn_added and plotly_cdn not in html_content:
        html_content = html_content.replace("</head>", f"{plotly_cdn}</head>")

    return html_content


class ReportExporter:
    def __init__(self):
        self.reports_dir = Path("data/reports")
        self.reports_dir.mkdir(parents=True, exist_ok=True)

    def export_pdf(self, sections: List[Dict], title: str = "Analysis Report", session_id: str = "") -> bytes:
        sections_html = ""
        for section in sections:
            section_title = section.get("title", "Section")
            content = section.get("content", "")
            content = self._embed_charts_as_images(content)
            sections_html += f'<section><h2>{section_title}</h2><div class="content">{content}</div></section>'

        full_html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
    body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 40px; color: #1e293b; line-height: 1.6; }}
    h1 {{ color: #1e3a5f; text-align: center; border-bottom: 2px solid #2563eb; padding-bottom: 10px; }}
    h2 {{ color: #2563eb; margin-top: 30px; }}
    .meta {{ text-align: center; color: #64748b; margin-bottom: 30px; }}
    section {{ margin-bottom: 20px; page-break-inside: avoid; }}
    .content {{ margin-left: 10px; }}
    table {{ border-collapse: collapse; width: 100%; margin: 10px 0; }}
    th, td {{ border: 1px solid #e2e8f0; padding: 8px; text-align: left; }}
    th {{ background-color: #f1f5f9; }}
    ul, ol {{ margin-left: 20px; }}
    strong, b {{ color: #0f172a; }}
    .chart-image {{ max-width: 100%; height: auto; margin: 15px 0; border: 1px solid #e2e8f0; }}
    .chart-placeholder {{ background: #f8fafc; border: 1px dashed #94a3b8; padding: 20px; text-align: center; margin: 10px 0; }}
</style>
</head>
<body>
<h1>{title}</h1>
<p class="meta">Generated: {datetime.now().strftime('%B %d, %Y')}</p>
{sections_html}
</body>
</html>"""

        from src.reporting.chart_converter import html_to_pdf_sync

        pdf_bytes = html_to_pdf_sync(full_html)
        if pdf_bytes:
            return pdf_bytes
        raise RuntimeError(
            "PDF generation failed. Ensure Playwright is installed with: pip install playwright && playwright install chromium"
        )

    def _embed_charts_as_images(self, content: str) -> str:
        iframe_pattern = r'<iframe[^>]*src="(/static/plots/|\.\./)([^"]+\.html)"[^>]*>.*?</iframe>'
        matches = re.findall(iframe_pattern, content, flags=re.DOTALL | re.IGNORECASE)

        for prefix, filename in matches:
            chart_path = Path("static/plots") / filename
            if chart_path.exists():
                try:
                    from src.reporting.chart_converter import convert_chart_to_png_sync, png_to_base64

                    png_bytes = convert_chart_to_png_sync(chart_path)
                    if png_bytes:
                        b64 = png_to_base64(png_bytes)
                        img_tag = f'<img src="data:image/png;base64,{b64}" class="chart-image" alt="{filename}"/>'
                        content = re.sub(
                            rf'<iframe[^>]*src="({re.escape(prefix)}{re.escape(filename)})"[^>]*>.*?</iframe>',
                            img_tag,
                            content,
                            flags=re.DOTALL | re.IGNORECASE,
                        )
                except Exception:
                    placeholder = f'<div class="chart-placeholder">[Chart: {filename}]</div>'
                    content = re.sub(
                        rf'<iframe[^>]*src="({re.escape(prefix)}{re.escape(filename)})"[^>]*>.*?</iframe>',
                        placeholder,
                        content,
                        flags=re.DOTALL | re.IGNORECASE,
                    )
        return content

    def export_docx(self, sections: List[Dict], title: str = "Analysis Report", session_id: str = "") -> bytes:
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            raise ImportError("python-docx is required for DOCX export. Install with: pip install python-docx")

        doc = Document()

        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        doc.add_paragraph()

        for section in sections:
            doc.add_heading(section.get("title", "Section"), level=1)
            content = section.get("content", "")
            if isinstance(content, str):
                clean_content = html_to_text(content)
                for para in clean_content.split("\n\n"):
                    if para.strip():
                        doc.add_paragraph(para)

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def export_xlsx(self, sections: List[Dict], data: Optional[Dict] = None, title: str = "Analysis Report") -> bytes:
        try:
            import openpyxl
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        except ImportError:
            raise ImportError("openpyxl is required for XLSX export. Install with: pip install openpyxl")

        wb = openpyxl.Workbook()

        summary_sheet = wb.active
        summary_sheet.title = "Summary"
        summary_sheet["A1"] = title
        summary_sheet["A1"].font = Font(size=16, bold=True)
        summary_sheet["A2"] = f"Generated: {datetime.now().strftime('%B %d, %Y')}"

        row = 4
        for section in sections:
            summary_sheet.cell(row=row, column=1, value=section.get("title", ""))
            summary_sheet.cell(row=row, column=1).font = Font(bold=True, size=12)
            row += 1
            content = section.get("content", "")
            if isinstance(content, str):
                for line in content.split("\n"):
                    if line.strip():
                        summary_sheet.cell(row=row, column=1, value=line.strip())
                        row += 1
            row += 1

        if data:
            if "dataframe" in data:
                import pandas as pd

                df = data["dataframe"]
                data_sheet = wb.create_sheet("Data")
                for c_idx, col in enumerate(df.columns, 1):
                    data_sheet.cell(row=1, column=c_idx, value=col)
                    data_sheet.cell(row=1, column=c_idx).font = Font(bold=True)
                for r_idx, row_data in enumerate(df.values, 2):
                    for c_idx, value in enumerate(row_data, 1):
                        data_sheet.cell(row=r_idx, column=c_idx, value=value)

        buffer = io.BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def export_pptx(self, sections: List[Dict], title: str = "Analysis Report", session_id: str = "") -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("python-pptx is required for PPTX export. Install with: pip install python-pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        subtitle_para.font.size = Pt(18)
        subtitle_para.alignment = PP_ALIGN.CENTER

        for section in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            section_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            section_title_frame = section_title_box.text_frame
            section_title_para = section_title_frame.paragraphs[0]
            section_title_para.text = section.get("title", "Section")
            section_title_para.font.size = Pt(32)
            section_title_para.font.bold = True

            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            content = section.get("content", "")
            if isinstance(content, str):
                paragraphs = content.split("\n\n")[:5]
                for i, para in enumerate(paragraphs):
                    if para.strip():
                        if i == 0:
                            p = content_frame.paragraphs[0]
                        else:
                            p = content_frame.add_paragraph()
                        p.text = para.strip()[:500]
                        p.font.size = Pt(16)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def export_html(self, sections: List[Dict], title: str = "Analysis Report", session_id: str = "") -> bytes:
        from html import escape

        html_parts = [
            "<!DOCTYPE html>",
            "<html lang='en'>",
            "<head>",
            f"<meta charset='UTF-8'><title>{escape(title)}</title>",
            "<style>",
            "body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; max-width: 900px; margin: 0 auto; padding: 40px; background: #0f172a; color: #e2e8f0; }",
            "h1 { color: #60a5fa; border-bottom: 2px solid #3b82f6; padding-bottom: 10px; }",
            "h2 { color: #93c5fd; margin-top: 30px; }",
            "p { line-height: 1.7; color: #cbd5e1; }",
            ".section { background: #1e293b; border-radius: 8px; padding: 20px; margin: 20px 0; border: 1px solid #334155; }",
            ".meta { color: #64748b; font-size: 14px; margin-bottom: 30px; }",
            "</style>",
            "</head>",
            "<body>",
            f"<h1>{escape(title)}</h1>",
            f"<p class='meta'>Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}</p>",
        ]

        for section in sections:
            section_title = escape(section.get("title", "Section"))
            content = section.get("content", "")
            if isinstance(content, str):
                content_html = content.replace("\n\n", "</p><p>").replace("\n", "<br>")
                content_html = f"<p>{content_html}</p>"
            else:
                content_html = str(content)

            html_parts.append(f"<div class='section'><h2>{section_title}</h2>{content_html}</div>")

        html_parts.extend(["</body>", "</html>"])
        return "\n".join(html_parts).encode("utf-8")

    def export_txt(self, sections: List[Dict], title: str = "Analysis Report", session_id: str = "") -> bytes:
        import re

        lines = [
            "=" * 60,
            title.upper(),
            "=" * 60,
            f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}",
            "",
        ]

        for section in sections:
            lines.append("-" * 40)
            lines.append(section.get("title", "Section").upper())
            lines.append("-" * 40)
            lines.append("")

            content = section.get("content", "")
            if isinstance(content, str):
                clean_content = re.sub(r"<[^>]+>", "", content)
                clean_content = clean_content.replace("&nbsp;", " ").replace("&amp;", "&")
                lines.append(clean_content)
            lines.append("")

        lines.append("=" * 60)
        lines.append("END OF REPORT")
        lines.append("=" * 60)

        return "\n".join(lines).encode("utf-8")

    def export_zip(self, sections: List[Dict], title: str = "Analysis Report", session_id: str = "") -> bytes:
        import zipfile
        import os

        buffer = io.BytesIO()

        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            if session_id:
                session_prefix = session_id[:8]
                report_dir = Path("data/reports")
                for report_file in report_dir.glob(f"*{session_prefix}*.html"):
                    if report_file.is_file():
                        zf.write(report_file, f"report/{report_file.name}")

            html_content = self.export_html(sections, title, session_id)
            zf.writestr("report/report_summary.html", html_content)

            txt_content = self.export_txt(sections, title, session_id)
            zf.writestr("report/report.txt", txt_content)

            artifacts_dir = Path("static/plots")
            if artifacts_dir.exists():
                for artifact_file in artifacts_dir.glob("*.html"):
                    if artifact_file.is_file():
                        zf.write(artifact_file, f"artifacts/{artifact_file.name}")
                for artifact_file in artifacts_dir.glob("*.png"):
                    if artifact_file.is_file():
                        zf.write(artifact_file, f"artifacts/{artifact_file.name}")

        buffer.seek(0)
        return buffer.getvalue()


_exporter_instance = None


def get_report_exporter() -> ReportExporter:
    global _exporter_instance
    if _exporter_instance is None:
        _exporter_instance = ReportExporter()
    return _exporter_instance
